import ast
import os
import logging
from copy import deepcopy
from dataclasses import dataclass, field, asdict
from typing import List, Optional
from pprint import pformat
from urllib.parse import urlparse

from celery import shared_task
from django.conf import settings
from pptx.util import Cm
import sendgrid
from sendgrid.helpers.mail import Mail, Email, Content, Attachment

import lookerapi as looker

import pickle
import os.path
from googleapiclient.discovery import build
from apiclient.http import MediaFileUpload
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request

# from .models import PresentationTemplate, DashboardTemplate, DashboardPresentation
from actions.actions import LookerInstance
from .models import PresentationPayload
from djax.utils import get_client


logger = logging.getLogger(__name__)

MAX_TILES = 4

GOOGLE_TOKEN = '/Users/looker/Documents/Github/djax/token.pickle'
GOOGLE_API_SCOPES = [
    'https://www.googleapis.com/auth/presentations.readonly',
    'https://www.googleapis.com/auth/drive',
    'https://www.googleapis.com/auth/presentations'
]

@dataclass
class LookerDashboardTile:
    id: int = None


@dataclass
class LookerDashboard:
    id: str = None
    tiles: List[LookerDashboardTile] = field(default_factory=list)


@dataclass
class LookerPresentationTile:
    title: str = None
    subtitle_text: str = None
    body: dict = None
    column: int = None
    row: int = None
    width: int = None
    height: int = None
    x: int = None
    y: int = None
    image_url: str = None
    image_binary: bytes = None
    tile_notes: str = None


@dataclass
class LookerPresentationSlide:
    title: str = None
    layout: str = None
    slide_type: str = None
    tiles: List[LookerPresentationTile] = field(default_factory=list)
    footnote: str = None
    slide_notes: str = None


@dataclass
class LookerPresentation:
    name: str = None
    dashboard_id: int = None
    template: str = None
    tiles: List[LookerPresentationTile] = field(default_factory=list)
    slides: List[LookerPresentationSlide] = field(default_factory=list)
    thumbnail_url: str = None
    thumbnail_binary: bytes = None


def convert_dashboard_tile_to_presentation_tile(element: dict) -> LookerPresentationTile:
    """Take the API response and convert into a LookerPresentationTile object"""
    return LookerPresentationTile(
        title=element['title'],
        subtitle_text=element['subtitle_text'],
        body=element['slide_params'],
        column=element['column'],
        row=element['row'],
        width=element['width'],
        height=element['height']
    )

def flush_tile_buffer_to_slide(tile_buffer: List[LookerPresentationTile]) -> Optional[LookerPresentationSlide]:
    """Take a list of KPI tile objects and convert into a LookerPresentationSlide object"""
    slide = LookerPresentationSlide(
        title='KPIs Slide',
        layout='Title Only',
        slide_type='Single Vis',
        tiles=deepcopy(tile_buffer),
        footnote='Link to dashboard TBD',
        slide_notes='Slide notes TBD'        
    )
    return slide

def convert_presentation_tile_to_slide(tile: LookerPresentationTile) -> Optional[LookerPresentationSlide]:
    """Take a single tile object and convert into a LookerPresentationSlide object"""
    slide = LookerPresentationSlide(
        title=tile.title,
        layout='Title Only',
        slide_type='Single Vis',
        tiles=[tile],
        footnote='Link to dashboard TBD',
        slide_notes='Slide notes TBD'        
    )
    return slide    

def get_image_file(client, tile: LookerPresentationTile, width=960, height=540):
    """By calling the Looker API, convert a tile object/definition into an actual .png image"""
    query_api = looker.QueryApi(client)
    element_request = {
        "body": tile.body,
        "result_format": 'png',
        "image_width": width,
        "image_height": height,
        "_preload_content": False,
    }

    image_file = os.path.join(settings.WORKING_DIR, f'{tile.title}.png')
    try:
        image = query_api.run_inline_query(**element_request)

        try:
            os.makedirs(os.path.dirname(image_file), exist_ok=True)
            with open(image_file, 'wb') as temp_image_file:
                temp_image_file.write(image.data)
        except Exception as e:
            logger.warning(f'Failed to save image for {tile.title}\n{e}')
            image_file = None

    except Exception as e:
        logger.error(f'API call failed for {tile.title}\n{pformat(element_request)}',
                     exc_info=True)
        image_file = None

    logger.info(f'Image file generated: {image_file}')
    return image_file

def add_image_to_slide(presentation_id, slide_id, slides_service, creds, image_file,
                       hyperlink_link=None, x=5, y=5, width=23):
    """
    Not yet working!!!
    Intended to:
    1. Take an image file for a given tile
    2. Upload it to Google Drive using MediaFileUpload()
    3. Add image to slide, using the Slide service batchUpdate() function
    """
    drive_service = build('drive', 'v3', credentials=creds, cache_discovery=False)

    file_metadata = {
        'name': f'SlideImage{slide_id}.png'
    }
    media = MediaFileUpload(image_file,
                            mimetype='image/png')
    file = drive_service.files().create(body=file_metadata,
                                        media_body=media
                                        ).execute()
    # img_url = f'https://drive.google.com/file/d/{file.get("id")}&access_token={creds.token}'
    img_url = f'https://drive.google.com/file/d/{file.get("id")}'
    logger.info(f'Image URL: {img_url}')
    logger.info(f'creds: {dir(creds)}')
    logger.info(f'token: {creds.token}')
    logger.info(f'files: {dir(file)}')
    logger.info(f'files: {file.keys()}')

    uploaded_file = drive_service.files().get(fileId=file.get("id"))
    logger.info(f'uploaded_file: {uploaded_file}')
    logger.info(f'uploaded_file: {dir(uploaded_file)}')
    logger.info(f'uploaded_file: {uploaded_file.to_json()}')

    requests = [
        {
            'createImage': {
                # 'url': 'https://www.betterbuys.com/wp-content/uploads/2015/07/Looker.png',
                'url': img_url,
                'elementProperties': {
                    'pageObjectId': slide_id,
                    'size': {
                        'width': {
                          'magnitude': 3000000,
                          'unit': 'EMU'
                        },
                        'height': {
                          'magnitude': 3000000,
                          'unit': 'EMU'
                        }
                      },
                      'transform': {
                        'scaleX': 0.6807,
                        'scaleY': 0.4585,
                        'translateX': 6583050,
                        'translateY': 1673950,
                        'unit': 'EMU'
                      }
                }
            }
        },
    ]
    slides_service.presentations().batchUpdate(body={'requests': requests},
            presentationId=presentation_id).execute()

    return file


# def add_footnote_to_slide(pptx, slide_no, hyperlink_text, hyperlink_link=None):
#     logger.info('dashboard_presentations.tasks.add_footnote_to_slide not yet implemented')
#     return

#     text_box = pptx.slides[slide_no].shapes.add_textbox(Cm(FOOT_X), Cm(FOOT_Y),
#                                                         width=Cm(FOOT_WIDTH), height=Cm(FOOT_HEIGHT))
#     p = text_box.text_frame.paragraphs[0]
#     run = p.add_run()
#     run.text = hyperlink_text
#     if hyperlink_link:
#         run.hyperlink.address = hyperlink_link
#     else:
#         run.hyperlink.address = hyperlink_text

def generate_google_slides_presentation(client, looker_presentation: LookerPresentation) -> Optional[str]:
    creds = None
    # The file token.pickle stores the user's access and refresh tokens, and is
    # created automatically when the authorization flow completes for the first
    # time.
    if os.path.exists(GOOGLE_TOKEN):
        with open(GOOGLE_TOKEN, 'rb') as token:
            creds = pickle.load(token)

    # If there are no (valid) credentials available, let the user log in.
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(
                'Google Slides/gslides_credentials.json', SCOPES)
            creds = flow.run_local_server()
        # Save the credentials for the next run
        with open('token.pickle', 'wb') as token:
            pickle.dump(creds, token)

    slides_service = build('slides', 'v1', credentials=creds, cache_discovery=False)

    body = {
        'title': looker_presentation.name
    }
    presentation = slides_service.presentations().create(body=body).execute()
    presentation_id = presentation.get('presentationId')
    logger.info(f'Created presentation with ID: {presentation_id}')

    for idx, slide in enumerate(looker_presentation.slides):
        slide_id = f'SlideNumber{idx}'
        title_id = f'Slide{idx}TitleField'
        requests = [
            {
                'createSlide': {
                    'objectId': slide_id,
                    'insertionIndex': idx,
                    'slideLayoutReference': {
                        'predefinedLayout': 'TITLE_ONLY'
                    },
                    'placeholderIdMappings': [
                        {
                            'layoutPlaceholder': {
                                'type': 'TITLE',
                                'index': 0,
                            },
                            'objectId': title_id
                        }
                    ]
                }
            },
            {
                'insertText': {
                    'objectId': title_id,
                    'text': slide.title
                }
            }
        ]

        body = {
            'requests': requests
        }
        response = slides_service.presentations().batchUpdate(
            presentationId=presentation_id,
            body=body).execute()

        if slide.tiles: #  == [5,4,3]
            image_file = get_image_file(client, slide.tiles[0])
            add_image_to_slide(presentation_id, slide_id, slides_service, creds, image_file)

        create_slide_response = response.get('replies')[0].get('createSlide')
        logger.info('Created slide with ID: {0}'.format(create_slide_response.get('objectId')))



@shared_task
# def generate_presentation_from_dashboard(instance: LookerInstanceRecord, payload: PresentationPayload) -> Optional[str]:
def generate_presentation_from_dashboard(payload: dict) -> Optional[str]:
    """
    Generates a report in PowerPoint format, using all looks in a Space.
    1.  Use template associated with this url, if available
    2.  Get dashboard definition from Dashboard API
    3.  Extract dashboard layout
    4.  Convert dashboard element objects to list of dicts – modify for content generation
    5.  Arrange elements for slides
    6.  Validate and prepare slidedeck
    7.  Loop list of tiles and add slides to LookerPresentation object
    8.  Convert LookerPresentation into Google Slides (or PowerPoint)
    9.  Register the generated slidedeck in the application database
    10. Email the report to recipients

    :param instance: LookerInstanceRecord containing details of the Looker instance hosting the dashboard
    :param payload: PresentationPayload with all data regarding the action request
    :return: Some kind of reference to the resulting presentation file
    """

    logger.info('generate_presentation_from_dashboard() called.')

    instance = LookerInstance(
        name=settings.INSTANCE,
        protocol=settings.INSTANCE_PROTOCOL,
        api_port=settings.INSTANCE_API_PORT,
        client_id=settings.USER_ID,
        client_secret=settings.USER_SECRET,
    )

    payload = PresentationPayload(**payload)

    parsed_url = urlparse(payload.url_with_params)
    dashboard_id = parsed_url.path.split('/')[2]

    # Get client to the Looker Instance
    try:
        client = get_client(instance)
    except Exception as e:
        logger.error('dashboard_presentations.tasks.py: Failed to get client', exc_info=True)
        return None

    # TODO: Handle filters
    # Get dashboard definition from Dashboard API
    dashboard_api = looker.DashboardApi(client)
    dashboard = dashboard_api.dashboard(dashboard_id)

    # Extract dashboard layout
    dashboard_layout = {}
    for component in dashboard.dashboard_layouts[0].dashboard_layout_components:
        dashboard_layout[component.dashboard_element_id] = {}
        dashboard_layout[component.dashboard_element_id]['column'] = component.column
        dashboard_layout[component.dashboard_element_id]['row'] = component.row
        dashboard_layout[component.dashboard_element_id]['width'] = component.width
        dashboard_layout[component.dashboard_element_id]['height'] = component.height

    ###########################################################
    # Convert dashboard element objects to list of dicts, then:
    # - Filter out vis elements
    # - 'Fix' the dictionary
    #    - have the query dict in a single location (varies between Tiles and Looks)
    #    – remove unwanted fields
    # - Enrich elements with layout information
    # - Sort elements based on layout information
    # - Group elements based on desired slide layout
    ############################################################
    elements = dashboard.dashboard_elements
    dashboard_tiles = [element.to_dict()
                           for element in elements
                           if element.type == 'vis']

    for element in dashboard_tiles:
        # Query dict in different locations depending on whether tile is also a Look
        if element['query']:
            element['slide_params'] = element['query']
        elif element['look']:
            element['slide_params'] = element['look']['query']
        else:
            logging.error(f'Could not find query for {element["id"]} {element["title"]}')
            return None

        # 1. Remove old client_id (a new client_id is generated by the server for each request)
        if 'client_id' in element['slide_params']:
            element['slide_params'].pop('client_id')

        # 2. The dynamic fields string has for some reason been split into individual chars,
        #    need to join them back together.
        if type(element['slide_params']['dynamic_fields']) == list:
            element['slide_params']['dynamic_fields'] = ''.join(element['slide_params']['dynamic_fields'])

        result_maker_vis_config = ast.literal_eval(element['result_maker']['vis_config'])
        element['slide_params']['vis_config'] = result_maker_vis_config

        # Enrich list of elements with layout info
        element['column'] = dashboard_layout[element['id']]['column']
        element['row'] = dashboard_layout[element['id']]['row']
        element['width'] = dashboard_layout[element['id']]['width']
        element['height'] = dashboard_layout[element['id']]['height']

    # Sort tiles to ensure the right order in the presentation
    dashboard_tiles = sorted(dashboard_tiles, key=lambda element: element['column'])
    dashboard_tiles = sorted(dashboard_tiles, key=lambda element: element['row'])

    # Start building LookerPresentation object
    presentation = LookerPresentation(
        name=dashboard.title
    )

    title_slide = LookerPresentationSlide(
        title=dashboard.title,
        slide_type='Title',
        layout='Title Only'
    )
    presentation.slides.append(title_slide)

    logger.info(f'dashboard_tile 0: {dashboard_tiles[0].keys()}')

    current_row = 0
    tile_buffer = []
    for raw_tile in dashboard_tiles:
        new_tile = convert_dashboard_tile_to_presentation_tile(raw_tile)

        if new_tile.body['vis_config']['type'] != 'single_value':
            if tile_buffer:
                new_slide = flush_tile_buffer_to_slide(tile_buffer)
                presentation.slides.append(new_slide)
                tile_buffer = []
            new_slide = convert_presentation_tile_to_slide(new_tile)
            presentation.slides.append(new_slide)

        else:
            if new_tile.row > current_row:
                if tile_buffer:
                    new_slide = flush_tile_buffer_to_slide(tile_buffer)
                    presentation.slides.append(new_slide)
                    tile_buffer = []

            current_row = new_tile.row
            tile_buffer.append(new_tile)
            if len(tile_buffer) == MAX_TILES:
                new_slide = flush_tile_buffer_to_slide(tile_buffer)
                presentation.slides.append(new_slide)
                tile_buffer = []

    # For dev purposes, dump the presentation contents to the log
    for idx, slide in enumerate(presentation.slides):
        tiles = ':'.join([tile.title for tile in slide.tiles ])
        logger.info(f'slide {idx}: {slide.title}: {tiles}')

    # Create presentation file
    generate_google_slides_presentation(client, presentation)

    # TODO: Register the generated slidedeck in the application database

    # TODO: Email the report to recipients

    return presentation
